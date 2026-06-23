"""Inject Marp speaker notes into the rendered PowerPoint file.

Marp CLI renders `--pptx` as image slides but does NOT copy the HTML-comment
speaker notes into the PowerPoint Notes pane. This script reads those notes
from the Marp `.md` (the single source of truth) and writes them into the
matching `.pptx` notes, slide by slide.

Run AFTER `marp ... --pptx`:
    python slides/apply_pptx_notes.py slides/module-1-power-bi.md slides/build/module-1-power-bi.pptx

Requires: python-pptx  (see pyproject.toml).
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation

# A comment is a Marp *directive* (not a note) when its first line looks like
# `key:` or `_key:` — e.g. `_class: lead`, `paginate: false`.
_DIRECTIVE = re.compile(r"^\s*_?[A-Za-z][\w-]*\s*:")
_COMMENT = re.compile(r"<!--(.*?)-->", re.DOTALL)


def notes_from_markdown(md_path: Path) -> list[str]:
    """Return one note string per slide, in slide order (empty if none)."""
    text = md_path.read_text(encoding="utf-8")
    parts = re.split(r"(?m)^---$", text)
    # parts[0] = preamble before first '---', parts[1] = front matter,
    # parts[2:] = slide bodies.
    slides = parts[2:]
    notes: list[str] = []
    for body in slides:
        slide_notes: list[str] = []
        for raw in _COMMENT.findall(body):
            content = raw.strip()
            first_line = content.splitlines()[0] if content else ""
            if _DIRECTIVE.match(first_line):
                continue  # it's a Marp directive, not a note
            # collapse the comment into a single tidy paragraph
            slide_notes.append(" ".join(content.split()))
        notes.append("\n\n".join(slide_notes))
    return notes


def main() -> None:
    if len(sys.argv) != 3:
        sys.exit("usage: python apply_pptx_notes.py <deck.md> <deck.pptx>")
    md_path, pptx_path = Path(sys.argv[1]), Path(sys.argv[2])
    notes = notes_from_markdown(md_path)
    prs = Presentation(str(pptx_path))

    n_pptx = len(prs.slides._sldIdLst)
    if len(notes) != n_pptx:
        print(
            f"WARNING: {len(notes)} markdown slides vs {n_pptx} pptx slides "
            "— applying by position."
        )

    applied = 0
    for slide, note in zip(prs.slides, notes):
        if not note:
            continue
        slide.notes_slide.notes_text_frame.text = note
        applied += 1

    prs.save(str(pptx_path))
    print(f"Applied notes to {applied} slide(s) in {pptx_path.name}")


if __name__ == "__main__":
    main()
